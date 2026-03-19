import os
import io
import base64
import PyPDF2
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from datetime import datetime
import concurrent.futures
import threading

# Thread-local storage for API clients to prevent rebuilding them per request
_thread_local = threading.local()

def get_thread_service(service_name, version, creds):
    key = f"{service_name}_{version}"
    if not hasattr(_thread_local, key):
        # static_discovery=False to avoid deprecated cache_discovery warnings but also prevent slowing down
        setattr(_thread_local, key, build(service_name, version, credentials=creds, cache_discovery=False))
    return getattr(_thread_local, key)

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build

SCOPES = [
    'https://www.googleapis.com/auth/gmail.readonly',
    'https://www.googleapis.com/auth/drive.readonly'
]

# ---------------------------------------------------------
# INTENT DICTIONARY & QUERY EXPANSION
# ---------------------------------------------------------
def expand_query(query):
    """Automatically expands query intent based on core terminology and synonyms."""
    query_lower = query.lower()
    expansions = {
        "resume": ["cv", "curriculum vitae", "biodata", "profile", "career summary", "experience"],
        "meeting notes": ["meeting summary", "discussion notes", "minutes of meeting", "mom", "sync notes"],
        "api": ["endpoint", "payload", "rest", "graphql", "interface", "authentication", "webservice"],
        "validation": ["error", "mismatch", "schema", "issue", "bug", "exception", "failed"],
        "report": ["summary", "analysis", "audit", "overview", "document", "whitepaper"],
        "policy": ["guidelines", "rules", "regulations", "compliance", "standard", "legal"],
        "mapping": ["integration", "matching", "alignment", "synchronization", "database schema"],
        "bmi": ["body mass index", "health report", "medical", "height weight", "fitness", "vitals", "biometric"],
        "certificate": ["certificate", "degree", "diploma", "qualification", "award", "transcript", "provisonal"],
        "kpi": ["key performance indicator", "performance", "metrics", "target", "goal", "quarterly", "report"],
        "gmeet": ["google meet", "meeting", "call", "video call", "conference", "discussion", "frd", "requirements"],
        "assignment": ["submission", "task", "homework", "project", "work", "marks", "grading"],
        "car": ["vehicle", "automobile", "insurance", "policy", "premium", "registration", "renewal", "rc", "chassis"],
        "policy": ["agreement", "contract", "insurance", "coverage", "terms", "policy document", "schedule", "premium"]
    }
    
    expanded_words = set(query_lower.split())
    for key, syns in expansions.items():
        if key in query_lower:
            expanded_words.update(syns)
            
    return " ".join(expanded_words)

# ---------------------------------------------------------
# INTELLIGENT SEMANTIC SCORING & RANKING
# ---------------------------------------------------------
class SemanticSearchEngine:
    def __init__(self, model_name='all-MiniLM-L6-v2'):
        self.model = SentenceTransformer(model_name)
        
    def search(self, query, documents, top_k=8, threshold=0.35):
        if not documents: return []
            
        print(f"DEBUG: Searching for '{query}' across {len(documents)} docs")
        expanded_query = expand_query(query)
        # Filter out common filler words that shouldn't trigger "Exact Match" boosts
        filler_words = {
            'help', 'find', 'out', 'the', 'for', 'document', 'with', 'mention', 'where', 
            'have', 'from', 'about', 'email', 'file', 'some', 'look', 'looking', 'where', 'search'
        }
        query_words = [w.lower() for w in query.split() if len(w) >= 3 and w.lower() not in filler_words]
        
        doc_keywords = ["resume", "cv", "report", "presentation", "policy", "document", "file", "profile", "biodata", "plan", "business"]
        has_doc_intent = any(kw in query.lower() for kw in doc_keywords)
        
        strings_to_encode = []
        mapping = []
        
        for idx, doc in enumerate(documents):
            # 1. Attachment Filename (Weight 3.0)
            if doc.get('attachment_name'):
                strings_to_encode.append(doc['attachment_name'])
                mapping.append((idx, 3.0, "Attachment Filename", doc['attachment_name']))
            
            # 2. File Title (Weight 2.8)
            if doc.get('title'):
                strings_to_encode.append(doc['title'])
                mapping.append((idx, 2.8, "File Name", doc['title']))

            # 3. Attachment Content (Weight 2.0)
            if doc.get('attachment_text'):
                strings_to_encode.append(doc['attachment_text'][:2000])
                mapping.append((idx, 2.0, "Attachment Content", doc['attachment_text']))
                
            # 4. Document/Email Content (Weight 1.5)
            content = doc.get('text', '')
            if content and len(content) > 5:
                strings_to_encode.append(content[:2000])
                reason = "Document Content" if doc.get('source') != "Gmail" else "Email Body"
                mapping.append((idx, 1.5, reason, content))

        if not strings_to_encode: return []

        embeddings = self.model.encode(strings_to_encode)
        query_emb = self.model.encode([expanded_query])
        
        sims = cosine_similarity(query_emb, embeddings)[0]
        
        doc_scores = {}
        for sim, (doc_idx, weight, reason, text) in zip(sims, mapping):
            if sim < threshold: continue
                
            weighted_sim = float(sim * weight)
            text_lower = text.lower()
            
            # STRICT MATCH LOGIC: Eliminate random drift
            match_multiplier = 1.0
            keywords_found = 0
            title_match = False
            
            for qw in query_words:
                if qw in text_lower:
                    keywords_found += 1
                    if reason in ["File Name", "Attachment Filename", "Email Subject"]:
                        title_match = True
            
            # If no keywords are found, we trust the semantic similarity less unless it's extremely high (>0.7)
            if keywords_found > 0:
                match_multiplier += (keywords_found * 0.2)
                if title_match: match_multiplier += 0.6
            elif sim < 0.6:
                # Penalize documents that don't have any of the keywords and have medium/low similarity
                match_multiplier = 0.5

            weighted_sim *= match_multiplier

            # Only show truly relevant items
            if weighted_sim < threshold: continue

            if doc_idx not in doc_scores or weighted_sim > doc_scores[doc_idx]['score']:
                # Snippet: Show clearer match context
                snippet_text = text[:800].replace('\n', ' ').strip()
                
                # Accuracy Labels
                if keywords_found > 1 and title_match:
                    label = "<b>🎯 High Accuracy Match</b>"
                elif sim > 0.6:
                    label = "<b>🔍 Strong Contextual Match</b>"
                else:
                    label = f"<b>Relevant Result</b> ({reason})"
                
                doc_scores[doc_idx] = {
                    'score': weighted_sim,
                    'explanation': f"{label}<br><div style='color: #444; font-size: 0.95rem; margin-top:5px;'><i>\"{snippet_text}...\"</i></div>",
                    'base_sim': float(sim)
                }

        results = []
        for doc_idx, data in doc_scores.items():
            res = documents[doc_idx].copy()
            res['score'] = data['score']
            res['explanation'] = data['explanation']
            results.append(res)
                
        # Rank Results
        results = sorted(results, key=lambda x: x['score'], reverse=True)
        return results[:top_k]

# ---------------------------------------------------------
# TEXT EXTRACTION FACTORY
# ---------------------------------------------------------
def extract_text_from_bytes(file_bytes, filename):
    """Zero-persistence text extraction for local files and fetched attachments"""
    text = ""
    ext = os.path.splitext(filename)[1].lower()
    file_obj = io.BytesIO(file_bytes)
    try:
        if ext == '.txt':
            text = file_bytes.decode('utf-8', errors='ignore')[:1500]
        elif ext == '.pdf':
            reader = PyPDF2.PdfReader(file_obj)
            if reader.pages: text = reader.pages[0].extract_text()[:1500]
        elif ext == '.docx':
            doc = docx.Document(file_obj)
            text = " ".join([p.text for p in doc.paragraphs if p.text.strip()])[:1500]
        elif ext == '.pptx':
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + " "
            text = text[:1500]
    except Exception:
        pass
    return text.strip()


# ---------------------------------------------------------
# SOURCE FETCHERS
# ---------------------------------------------------------
def fetch_local_files(directory_path, max_files=50):
    docs = []
    if not os.path.exists(directory_path): return docs
    
    for root, dirs, files in os.walk(directory_path):
        for file in files:
            if len(docs) >= max_files: break
            file_path = os.path.join(root, file)
            ext = os.path.splitext(file)[1].lower()
            if ext not in ['.txt', '.pdf', '.docx', '.pptx']: continue
            
            try:
                mod_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%Y-%m-%d')
                with open(file_path, 'rb') as f:
                    content = extract_text_from_bytes(f.read(), file)
                
                if content:
                    docs.append({
                        "id": file_path, "title": file, "source": "Local File",
                        "sender": directory_path, "date": mod_time, "text": content,
                        "attachment_name": None, "attachment_text": "", "type": "document"
                    })
            except Exception: pass
    return docs

from google.auth.exceptions import RefreshError

def get_google_credentials():
    creds = None
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    
    if not creds or not creds.valid:
        try:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                raise RefreshError("No valid credentials or refresh token")
        except (RefreshError, Exception):
            if not os.path.exists('credentials.json'): 
                return None
            
            # Local flow - Streamlit Cloud will need st.secrets update
            try:
                flow = InstalledAppFlow.from_client_secrets_file('credentials.json', SCOPES)
                creds = flow.run_local_server(port=0, open_browser=False)
            except Exception as e:
                print(f"Authentication Error: {e}")
                return None
        
        # Save token
        if creds:
            with open('token.json', 'w') as token:
                token.write(creds.to_json())
                
    return creds

def fetch_gmail(query_text, max_results=15):
    """Parallelized Gmail search with strict filtering."""
    creds = get_google_credentials()
    if not creds: return []
    try:
        service = build('gmail', 'v1', credentials=creds)
        expanded = expand_query(query_text)
        doc_keywords = ["resume", "cv", "report", "presentation", "policy", "document", "file", "profile", "biodata", "plan"]
        has_doc_intent = any(kw in query_text.lower() for kw in doc_keywords)
        
        base_query = " OR ".join(expanded.split())
        query_words = [w.lower() for w in query_text.split() if len(w) >= 2]
        filter_str = "category:primary -category:promotions -category:social -category:updates -category:forums -in:spam -in:trash"
        marketing_exclude = "-unsubscribe -sale -offer -discount -deal -\"limited time\" -newsletter"
        
        # Build a service for the main list call
        service = build('gmail', 'v1', credentials=creds, cache_discovery=False)
        
        gmail_query = f"({base_query}) {filter_str}" # Broadened query
        if has_doc_intent: gmail_query = f"({base_query}) has:attachment {filter_str}"
        
        # Pull 60 messages for maximum contextual coverage
        results = service.users().messages().list(userId='me', maxResults=60, q=gmail_query).execute()
        messages = results.get('messages', [])
        
        def process_msg(msg):
            try:
                # Thread-safety: Build service inside the worker (optimized with thread local)
                t_service = get_thread_service('gmail', 'v1', creds)
                msg_data = t_service.users().messages().get(userId='me', id=msg['id'], format='full').execute()
                headers = msg_data.get('payload', {}).get('headers', [])
                subject, sender, date = "No Subject", "Unknown Sender", "Unknown Date"
                is_bulk = False
                for header in headers:
                    name = header['name'].lower()
                    val = header['value']
                    if name == 'subject': subject = val
                    elif name == 'from': sender = val
                    elif name == 'date': date = val
                    elif name == 'list-unsubscribe' or (name == 'precedence' and val.lower() in ['bulk', 'list', 'junk']): is_bulk = True
                
                if is_bulk or any(ms in sender.lower() for ms in ["noreply@", "no-reply@", "newsletter@"]): return None

                body_text = ""
                attachment_name, attachment_text = None, ""
                payload = msg_data.get('payload', {})
                parts = [payload] if 'parts' not in payload else payload['parts']
                
                def parse_parts(parts_list):
                    nonlocal body_text, attachment_name, attachment_text
                    for part in parts_list:
                        mime, filename, att_id = part.get('mimeType'), part.get('filename'), part.get('body', {}).get('attachmentId')
                        if filename and att_id:
                            attachment_name = filename
                            try:
                                att = t_service.users().messages().attachments().get(userId='me', messageId=msg['id'], id=att_id).execute()
                                attachment_text = extract_text_from_bytes(base64.urlsafe_b64decode(att.get('data')), filename)
                            except: pass
                        elif mime == 'text/plain' and part.get('body', {}).get('data'):
                            body_text += base64.urlsafe_b64decode(part['body']['data']).decode('utf-8', errors='ignore')
                        elif mime == 'text/html' and part.get('body', {}).get('data'):
                            html = base64.urlsafe_b64decode(part['body']['data']).decode('utf-8', errors='ignore')
                            body_text += BeautifulSoup(html, 'html.parser').get_text()
                        elif 'parts' in part: parse_parts(part['parts'])
                
                parse_parts(parts)
                return {
                    "id": msg['id'], "title": subject, "source": "Gmail", "sender": sender[:60], "date": date[:16], 
                    "text": body_text[:1000].strip(), "attachment_name": attachment_name, "attachment_text": attachment_text, "type": "email",
                    "url": f"https://mail.google.com/mail/u/0/#inbox/{msg['id']}"
                }
            except: return None

        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
            docs = list(filter(None, executor.map(process_msg, messages)))
        return docs
    except Exception as e:
        print(f"Gmail Error: {e}")
        return []

def fetch_google_drive(query_text, max_results=20):
    """Parallelized Google Drive search with deep content extraction."""
    creds = get_google_credentials()
    if not creds: return []
    try:
        service = build('drive', 'v3', credentials=creds)
        expanded = expand_query(query_text)
        # Filter out common filler words for the Drive API query to avoid invalid/long queries
        stop_words = {'help', 'find', 'out', 'the', 'for', 'document', 'with', 'mention', 'where', 'have', 'from', 'about'}
        clean_words = [w for w in expanded.split() if w.lower() not in stop_words and len(w) >= 2]
        
        # Search BOTH name and content for much better recall
        drive_query = " or ".join([f"(name contains '{word}' or fullText contains '{word}')" for word in clean_words[:10]])
        
        # Include images and screenshots in the search as certificates/assignments are often scans
        image_query = "mimeType='image/jpeg' or mimeType='image/png' or mimeType='application/pdf' or mimeType='application/vnd.google-apps.document'"
        if not drive_query: drive_query = image_query
            
        # BROAD RETRIEVAL: Pull top 60 candidates. We rely on the LOCAL semantic engine to rank them.
        results = service.files().list(pageSize=60, fields="files(id, name, mimeType, description, modifiedTime, owners)", q=f"trashed=false and ({drive_query})").execute()
        items = results.get('files', [])
        
        def process_file(item):
            try:
                # Thread-safety: Build service inside the worker (optimized with thread local)
                t_service = get_thread_service('drive', 'v3', creds)
                file_id, title, mime = item['id'], item.get('name', 'Unknown'), item.get('mimeType')
                date = item.get('modifiedTime', 'Unknown Date')
                owners = item.get('owners', [])
                sender = owners[0].get('displayName', 'My Drive') if owners else 'My Drive'
                content = item.get('description', '')
                try:
                    if 'google-apps.document' in mime:
                        content = t_service.files().export(fileId=file_id, mimeType='text/plain').execute().decode('utf-8')[:2000]
                    elif mime == 'application/pdf':
                        content = extract_text_from_bytes(t_service.files().get_media(fileId=file_id).execute(), title)
                except: pass
                return {"id": file_id, "title": title, "source": "Google Drive", "sender": sender, "date": date[:10], "text": content.strip(), "attachment_name": None, "attachment_text": "", "type": "document", "url": f"https://drive.google.com/file/d/{file_id}/view" }
            except: return None

        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
            docs = list(filter(None, executor.map(process_file, items)))
        return docs
    except Exception as e:
        print(f"Drive Error: {e}")
        return []
